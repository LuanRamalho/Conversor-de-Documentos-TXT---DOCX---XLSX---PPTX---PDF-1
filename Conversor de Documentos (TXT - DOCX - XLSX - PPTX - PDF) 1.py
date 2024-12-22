import os
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from docx import Document
from pptx import Presentation
import pandas as pd
from reportlab.pdfgen import canvas

# Função para converter TXT para PDF
def convert_txt_to_pdf(input_file, output_file):
    try:
        with open(input_file, 'r', encoding='utf-8') as file:
            text = file.readlines()
        pdf = canvas.Canvas(output_file)
        for i, line in enumerate(text):
            pdf.drawString(50, 800 - i * 15, line.strip())
        pdf.save()
        return True
    except Exception as e:
        messagebox.showerror("Erro", f"Erro ao converter TXT: {e}")
        return False

# Função para converter DOCX para PDF
def convert_docx_to_pdf(input_file, output_file):
    try:
        doc = Document(input_file)
        pdf = canvas.Canvas(output_file)
        y = 800
        for paragraph in doc.paragraphs:
            pdf.drawString(50, y, paragraph.text)
            y -= 15
            if y <= 50:
                pdf.showPage()
                y = 800
        pdf.save()
        return True
    except Exception as e:
        messagebox.showerror("Erro", f"Erro ao converter DOCX: {e}")
        return False

# Função para converter XLSX para PDF
def convert_xlsx_to_pdf(input_file, output_file):
    try:
        df = pd.read_excel(input_file)
        pdf = canvas.Canvas(output_file)
        y = 800
        for column in df.columns:
            pdf.drawString(50, y, column)
            y -= 15
            for value in df[column]:
                pdf.drawString(70, y, str(value))
                y -= 15
                if y <= 50:
                    pdf.showPage()
                    y = 800
        pdf.save()
        return True
    except Exception as e:
        messagebox.showerror("Erro", f"Erro ao converter XLSX: {e}")
        return False

# Função para converter PPTX para PDF
def convert_pptx_to_pdf(input_file, output_file):
    try:
        ppt = Presentation(input_file)
        pdf = canvas.Canvas(output_file)
        for slide in ppt.slides:
            y = 800
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        pdf.drawString(50, y, paragraph.text)
                        y -= 15
                        if y <= 50:
                            pdf.showPage()
                            y = 800
        pdf.save()
        return True
    except Exception as e:
        messagebox.showerror("Erro", f"Erro ao converter PPTX: {e}")
        return False

# Função para realizar a conversão
def convert_file():
    input_file = file_path.get()
    if not input_file:
        messagebox.showwarning("Aviso", "Por favor, selecione um arquivo.")
        return
    
    file_extension = input_file.split('.')[-1].lower()
    output_file = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
    if not output_file:
        return

    if file_extension == "txt":
        success = convert_txt_to_pdf(input_file, output_file)
    elif file_extension == "docx":
        success = convert_docx_to_pdf(input_file, output_file)
    elif file_extension == "xlsx":
        success = convert_xlsx_to_pdf(input_file, output_file)
    elif file_extension == "pptx":
        success = convert_pptx_to_pdf(input_file, output_file)
    else:
        messagebox.showerror("Erro", "Formato de arquivo não suportado.")
        return

    if success:
        messagebox.showinfo("Sucesso", f"Arquivo convertido com sucesso para {output_file}.")

# Função para selecionar o arquivo
def select_file():
    filetypes = [("Arquivos suportados", "*.txt *.docx *.xlsx *.pptx")]
    file = filedialog.askopenfilename(filetypes=filetypes)
    if file:
        file_path.set(file)

# Configuração da interface gráfica
root = tk.Tk()
root.title("Conversor de Arquivos")
root.geometry("500x300")
root.resizable(False, False)

file_path = tk.StringVar()

# Estilo
style = ttk.Style()
style.configure("TCombobox", font=("Arial", 12))

frame = tk.Frame(root, bg="#FF80F5")
frame.place(relwidth=1, relheight=1)

title_label = tk.Label(frame, text="Conversor de Arquivos para PDF", font=("Arial", 16, "bold"), bg="#4CAF50", fg="white")
title_label.pack(fill="x", pady=10)

file_label = tk.Label(frame, text="Selecione o arquivo:", font=("Arial",12), bg="#FF80F5")
file_label.pack(pady=5)

file_entry = tk.Entry(frame, textvariable=file_path, width=50, state="readonly", font=("Arial",12), bg="#FF80F5")
file_entry.pack(pady=5)

file_button = tk.Button(frame, text="Escolher Arquivo", command=select_file, font=("Arial",12,"bold"), bg="#020E80", fg="#ffffff")
file_button.pack(pady=5)

convert_button = tk.Button(frame, text="Converter para PDF", command=convert_file, font=("Arial",12,"bold"), bg="#036701", fg="#ffffff")
convert_button.pack(pady=20)

root.mainloop()
